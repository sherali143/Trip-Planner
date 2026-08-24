"""
WHAT THIS FILE DOES
===================
Generates CMP7200_Viva_Presentation.pptx — the deck for the viva voce, which is
20% of the module.

    python -m report.build.make_viva_deck

Why it is generated rather than written by hand
-----------------------------------------------
Every number on a slide is read from evaluation/measured.py, the same accessor the
dissertation and the project overview use. A deck typed by hand is a third place
for the same figure to be wrong in, and the one place nobody re-checks before
presenting.

Why there is so little text on each slide
----------------------------------------
The viva marking criteria penalise it directly. The 30-39% band for the visual
criterion reads "Slides lack information on key aspects of project, too
text-heavy", and the bands above reward "encapsulation of key points". So each
slide carries one idea: a headline, a figure or a small number of short lines, and
nothing that the presenter would end up reading aloud.

The speaker notes carry the detail instead. They are not projected, so they cost
nothing on the slide and give the presenter the argument, the caveat and the
figure's provenance in the place where those are actually needed.

Figures are the generated ones under report/figures/, so the deck and the
dissertation show the reader the same charts.
"""

from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(
    os.path.abspath(__file__)))))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from evaluation import measured

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIGURES = os.path.join(ROOT, "report", "figures")
OUTPUT = os.path.join(ROOT, "CMP7200_Viva_Presentation.pptx")

STUDENT = "25182589"          # anonymous marking: student number only
MODULE = "CMP7200 — Individual Master's Project"

INK = RGBColor(0x16, 0x1A, 0x23)
SOFT = RGBColor(0x54, 0x5C, 0x6E)
FAINT = RGBColor(0x8A, 0x91, 0xA0)
ACCENT = RGBColor(0x43, 0x38, 0xCA)
TEAL = RGBColor(0x0E, 0xA5, 0xA4)
WARN = RGBColor(0xA8, 0x5B, 0x00)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1B, 0x4B)

W, H = Inches(13.333), Inches(7.5)      # 16:9


# ---------------------------------------------------------------------------
# Slide furniture
# ---------------------------------------------------------------------------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """
    Add a text box. `runs` is a list of (text, size_pt, bold, colour) tuples,
    one per paragraph, so a slide's type hierarchy is declared where it is used.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    for index, (text, size, bold, colour) in enumerate(runs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Calibri"
    return box


def _rule(slide, left, top, width, colour=ACCENT, height=Emu(34000)):
    bar = slide.shapes.add_shape(1, left, top, width, height)   # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = " ".join(text.split())


def _heading(slide, eyebrow, title):
    """The standard slide head: a small label, a large title, a short rule."""
    _text(slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.4),
          [(eyebrow.upper(), 12, True, ACCENT)])
    _text(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.9),
          [(title, 32, True, INK)])
    _rule(slide, Inches(0.9), Inches(1.72), Inches(1.1))


def _figure(slide, name, top=Inches(2.0), height=Inches(4.6)):
    """Place a generated figure, centred, or say plainly that it is missing."""
    for folder in ("results", "diagrams"):
        path = os.path.join(FIGURES, folder, name)
        if os.path.exists(path):
            picture = slide.shapes.add_picture(path, Inches(0), top,
                                               height=height)
            picture.left = int((W - picture.width) / 2)
            return picture
    _text(slide, Inches(0.9), top, Inches(11.5), Inches(0.5),
          [(f"[figure {name} not generated — run make_charts and make_diagrams]",
            14, False, WARN)])
    return None


def _stat_row(slide, stats, top=Inches(2.3), colour=ACCENT):
    """
    Evenly spaced big numbers. This is the deck's main device: a figure a marker
    can read from the back of the room, with three words under it.
    """
    count = len(stats)
    gutter = Inches(0.9)
    span = (W - 2 * gutter) / count
    for index, (value, label) in enumerate(stats):
        left = int(gutter + index * span)
        _text(slide, Emu(left), top, Emu(int(span)), Inches(1.3),
              [(value, 60, True, colour)], align=PP_ALIGN.CENTER)
        _text(slide, Emu(left), top + Inches(1.25), Emu(int(span)), Inches(0.8),
              [(label, 14, False, SOFT)], align=PP_ALIGN.CENTER)


def _table(slide, headers, rows, left=Inches(0.9), top=Inches(2.1),
           width=Inches(11.5), col_widths=None, font=12):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top,
                                   width, Inches(0.4 * (len(rows) + 1)))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for index, share in enumerate(col_widths):
            table.columns[index].width = Emu(int(width * share / total))
    def style(cell, value, bold, colour, fill):
        # Assigning "" leaves the paragraph with no runs, so styling it raises.
        # An empty header is legitimate — the approach table's first column is
        # just the letter — so write a space and style that.
        cell.text = str(value) if str(value) else " "
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(font)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

    for index, head in enumerate(headers):
        style(table.cell(0, index), head, True, PAPER, ACCENT)
    for r, row in enumerate(rows, start=1):
        shade = PAPER if r % 2 else RGBColor(0xF7, 0xF8, 0xFC)
        for c, value in enumerate(row):
            style(table.cell(r, c), value, False, INK, shade)
    return table


# ---------------------------------------------------------------------------
# The deck
# ---------------------------------------------------------------------------

def build() -> str:
    arms = measured.results()["arms"]
    gains = measured.results().get("improvements", {})
    coverage = measured.coverage()
    protocol = measured.protocol_summary()
    gate = measured.gate_agreement()
    cache = measured.api_cache_stats()
    quota = measured.api_quota()["apis"]
    code = measured.code_stats()
    tests = measured.test_count()["collected"]

    def calls(code_letter):
        return arms[code_letter]["avg_llm_calls"]

    def tokens(code_letter):
        return arms[code_letter]["avg_total_tokens"]

    def grounded(code_letter):
        return arms[code_letter]["avg_prices_grounded_pct"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------------- 1 title
    slide = _blank(prs)
    _fill(slide, DARK)
    _text(slide, Inches(1.1), Inches(1.9), Inches(11.2), Inches(0.5),
          [("MSc DISSERTATION · VIVA VOCE", 14, True,
            RGBColor(0x9C, 0xA3, 0xF5))])
    _text(slide, Inches(1.1), Inches(2.4), Inches(11.2), Inches(1.9),
          [("Does a language model earn its cost", 40, True, PAPER),
           ("in a multi-agent travel planner?", 40, True,
            RGBColor(0x6E, 0xE7, 0xDF))], spacing=1.1)
    _rule(slide, Inches(1.1), Inches(4.5), Inches(1.6), TEAL)
    _text(slide, Inches(1.1), Inches(4.9), Inches(11.2), Inches(1.4),
          [("Four architectures measured on the same task, "
            "with retrieval as the only variable.", 17, False,
            RGBColor(0xC7, 0xD2, 0xFE)),
           ("", 8, False, PAPER),
           (f"Student {STUDENT}   ·   {MODULE}", 13, False,
            RGBColor(0x9C, 0xA3, 0xF5))])
    _notes(slide, """
        Fifteen minutes. The argument is one sentence: a model is worth its cost
        where a step needs judgement, and costs without adding anything where it
        does not. I built four architectures that differ only in how data is
        retrieved, measured all four, and the retrieval layer is where the money
        goes.
    """)

    # ---------------------------------------------------------------- 2 problem
    slide = _blank(prs)
    _heading(slide, "The problem", "Agents are given tools by default")
    _stat_row(slide, [
        (f"{calls('B'):.0f}", "model requests per trip\nwhen six agents each decide"),
        (f"{calls('D'):.0f}", "model requests per trip\nwhen plain Python fetches"),
        (f"{tokens('B') / tokens('D'):.0f}x", "the tokens, for the same\nfinished itinerary"),
    ], colour=ACCENT)
    _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
          [("Fetching a fare for a known route on a known date needs no judgement.",
            20, True, INK),
           ("A reasoning loop is charged for it anyway — every step re-sends every "
            "tool schema.", 16, False, SOFT)], spacing=1.25)
    _notes(slide, f"""
        The literature treats tool-use as the thing that makes agents useful. Nobody
        asks which steps deserve one. Interpreting "leaving on the fifteenth for two
        adults" is a judgement. Calling a flight API with a date is not. The naive
        arm spent {calls('B'):.1f} model requests per trip against
        {calls('D'):.1f} for the shipped design, and the tokens differ by about
        {tokens('B') / tokens('D'):.0f} times, because each loop step re-sends the
        whole tool schema set.
    """)

    # ---------------------------------------------------------------- 3 question
    slide = _blank(prs)
    _fill(slide, RGBColor(0xF7, 0xF8, 0xFC))
    _text(slide, Inches(1.4), Inches(2.3), Inches(10.5), Inches(2.6),
          [("Where does a language model", 34, False, SOFT),
           ("earn its cost,", 46, True, ACCENT),
           ("and where is it just expensive?", 34, False, SOFT)], spacing=1.15)
    _rule(slide, Inches(1.4), Inches(5.2), Inches(2.0), TEAL)
    _text(slide, Inches(1.4), Inches(5.5), Inches(10.5), Inches(0.8),
          [("Design Science Research · four-arm ablation · retrieval is the "
            "only variable", 15, False, FAINT)])
    _notes(slide, """
        Everything else is held constant: same task, same prompts where they
        overlap, same tool server, same message protocol, same model. Only the
        retrieval layer changes between arms. That is what makes the comparison
        attributable.
    """)

    # ---------------------------------------------------------------- 4 approaches
    slide = _blank(prs)
    _heading(slide, "The four approaches", "One variable: who fetches the data")
    _table(slide,
           ["", "Architecture", "Who retrieves", "Why it is in the design"],
           [["A", "One model call, no tools", "nothing",
             "Floor. Shows what fluency alone produces"],
            ["B", "Six agents, as first built", "each agent decides",
             "The proposal, untuned. The honest baseline"],
            ["C", "Six agents, tuned", "each agent decides",
             "Separates tuning from architecture"],
            ["D", "Three agents", "plain Python",
             "What ships. The claim under test"]],
           col_widths=[0.5, 3.0, 2.2, 5.0], font=13, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0),
          [("C exists so nobody can say D only won because B was badly configured.",
            18, True, ACCENT)])
    _notes(slide, f"""
        C is the arm that makes the result defensible. Tuning alone cut tokens by
        {gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}% without changing the
        architecture — so a comparison of B against D would have conflated two
        different effects. A is the control: it produces a fluent itinerary from
        nothing, and its price groundedness of {grounded('A'):.1f}% is the number
        that shows why retrieval matters at all.
    """)

    # ---------------------------------------------------------------- 5 architecture
    slide = _blank(prs)
    _heading(slide, "Architecture", "Four layers; only layer three changes")
    _figure(slide, "architecture.png", top=Inches(2.0), height=Inches(4.7))
    _notes(slide, """
        Layer one takes the request, two turns free text into typed fields, three
        retrieves, four assembles. The research question lives entirely in layer
        three. Note the honest detail: the shipped path imports the tool functions
        in process. The JSON-RPC transport is exercised by the six-agent arms, and
        Section 7.2 of the dissertation says so rather than implying otherwise.
    """)

    # ---------------------------------------------------------------- 6 agents
    slide = _blank(prs)
    _heading(slide, "The agents", "Three use the model. Four steps do not")
    _table(slide,
           ["Step", "Who", "Model?", "Why"],
           [["1", "Conversational agent", "yes",
             "Reads a request written however a person writes it"],
            ["2", "Preferences extractor", "yes",
             "Free text to typed fields, then the budget is checked"],
            ["3", "Flights · hotels · attractions · restaurants", "no",
             "One correct call each. An if statement decides, not a model"],
            ["4", "Itinerary coordinator", "yes",
             "Arranges retrieved options into a sensible day-by-day plan"]],
           col_widths=[0.5, 3.4, 0.9, 5.9], font=13, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.1),
          [("The test applied to every step: does this need a judgement that "
            "cannot be written as code?", 18, True, INK),
           ("Two steps do. Four do not. That is the whole design.", 16, False,
            SOFT)], spacing=1.25)
    _notes(slide, """
        Six agents became three because the instrumentation showed most requests
        were spent on deterministic retrieval. The three that remain each fail the
        "write it as code" test — you cannot write a rule that interprets arbitrary
        phrasing, and you cannot write a rule that sequences a pleasant day.
    """)

    # ---------------------------------------------------------------- 7 method
    slide = _blank(prs)
    _heading(slide, "How it was measured", "Repeats, intervals, and no credentials")
    _stat_row(slide, [
        (f"{coverage['repeats_per_arm']}", "runs per arm, so every\nfigure has an interval"),
        (f"{coverage['scenarios_measured']}/{coverage['scenarios_designed']}",
         "scenarios for the cost\ncomparison — quota-bound"),
        (f"{cache['entries']}", "recorded API responses,\ncommitted and replayable"),
        (f"{tests}", "tests, no keys\nand no network"),
    ], colour=ACCENT)
    _text(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2),
          [("Model requests counted from provider callbacks, never by hand.",
            17, False, SOFT),
           ("Coverage is stated on every chart rather than smoothed over.",
            17, False, SOFT)], spacing=1.3)
    _notes(slide, f"""
        The honest weakness first: {coverage['scenarios_measured']} of
        {coverage['scenarios_designed']} scenarios for the cost comparison, because
        the free tiers allow thirty flight and fifty hotel searches a month. What
        that bought instead: two questions that need no quota — protocol
        conformance and the feasibility gate — evaluated across all twenty, and
        both produced negative findings. Never hand-count model requests: a
        reasoning loop issues far more than there are tasks.
    """)

    # ---------------------------------------------------------------- 8 cost
    slide = _blank(prs)
    _heading(slide, "Result · cost", "The retrieval layer is where the money goes")
    _figure(slide, "efficiency.png", top=Inches(1.95), height=Inches(4.4))
    _text(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
          [(f"B {tokens('B'):,.0f} tokens  →  D {tokens('D'):,.0f}.  "
            f"Same task, same data, same model.", 16, True, INK)])
    _notes(slide, f"""
        Naive six agents: {tokens('B'):,.0f} tokens and {calls('B'):.1f} model
        requests per trip. Three agents with Python retrieval:
        {tokens('D'):,.0f} tokens and {calls('D'):.1f} requests. The intervals do
        not overlap, which is why this is reported as a difference rather than as
        noise. Latency follows the same shape.
    """)

    # ---------------------------------------------------------------- 9 tuning
    slide = _blank(prs)
    _heading(slide, "Result · tuning vs architecture",
             "Configuration is not the same as design")
    _figure(slide, "tuning_effect.png", top=Inches(1.95), height=Inches(4.4))
    _text(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
          [(f"Tuning alone: {gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}% "
            f"fewer tokens. Changing the architecture: "
            f"{gains.get('D_vs_B', {}).get('tokens_pct', 0):.0f}%.",
            16, True, INK)])
    _notes(slide, f"""
        This slide is the answer to the obvious objection. Cutting each agent's
        tool list from twenty slots to four, and its step budget down, took
        {gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}% of the tokens out
        without touching the architecture. So most of B's cost was configuration.
        The remainder — the part only the architecture change reaches — is the
        contribution being claimed.
    """)

    # ---------------------------------------------------------------- 10 grounded
    slide = _blank(prs)
    _heading(slide, "Result · groundedness", "Cheaper, and no less real")
    _figure(slide, "groundedness.png", top=Inches(1.95), height=Inches(4.4))
    _text(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
          [(f"Prices traceable to retrieved data — A {grounded('A'):.1f}%  ·  "
            f"B {grounded('B'):.1f}%  ·  C {grounded('C'):.1f}%  ·  "
            f"D {grounded('D'):.1f}%", 16, True, INK)])
    _notes(slide, f"""
        The cost saving would be worthless if the plan got vaguer. A is the control
        and scores {grounded('A'):.1f}% — a fluent itinerary invented from nothing,
        which is exactly the failure this metric exists to detect. It caught a real
        one: the first end-to-end run produced a confident plan while every tool
        call was silently failing.
    """)

    # ---------------------------------------------------------------- 11 negative
    slide = _blank(prs)
    _heading(slide, "Findings that were not flattering",
             "Two things I built do not do what I claimed")
    _stat_row(slide, [
        (f"{protocol['passed']}/{protocol['total_checks']}",
         "protocol conformance checks pass.\nPriority is declared, never honoured"),
        (f"{gate['cohens_kappa']:.3f}",
         "Cohen's kappa on the budget gate.\nModerate, not good"),
    ], colour=WARN)
    _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
          [("Both are reported rather than quietly fixed.", 20, True, INK),
           (f"The gate misses {gate['false_negative']} of "
            f"{gate['true_positive'] + gate['false_negative']} impossible budgets "
            f"— recall {gate['recall']:.1f}, precision {gate['precision']:.1f}.",
            16, False, SOFT)], spacing=1.25)
    _notes(slide, f"""
        {protocol['failed']} of {protocol['total_checks']} conformance checks fail:
        message priority is declared and never honoured, and inbound permissions
        are never enforced. On the gate, precision is {gate['precision']:.1f} and
        recall {gate['recall']:.1f} — it never wrongly refuses a workable trip, and
        it misses half the impossible ones. Kappa {gate['cohens_kappa']:.3f} is
        moderate agreement. Saying so is the point; an evaluation that only
        confirms its author is not an evaluation.
    """)

    # ---------------------------------------------------------------- 12 risk
    slide = _blank(prs)
    _heading(slide, "Risk", "The binding constraint was never difficulty")
    flight = next(v for v in quota.values() if "flight" in v["name"])
    hotel = next(v for v in quota.values() if "hotel" in v["name"])
    _stat_row(slide, [
        (f"{int(flight['limit'])}", "flight searches a month.\nOne search costs two"),
        (f"{int(hotel['limit'])}", "hotel searches a month.\nFree tier, no top-up"),
        (f"{cache['entries']}", "responses recorded once,\nreplayed for ever"),
    ], colour=WARN)
    _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
          [("Exhausting the quota is not a delay. It is a month-long stop.",
            20, True, INK),
           ("Record-and-replay plus a hard live-call ceiling protected "
            "reproducibility — not breadth.", 16, False, SOFT)], spacing=1.25)
    _notes(slide, f"""
        Both allowances were fully spent during development. The controls were built
        before the arms were, and they worked: {cache['entries']} committed
        responses mean a marker with no API keys can reproduce every number. What
        they could not do is create requests that were never made, so coverage of
        the cost comparison stayed at {coverage['scenarios_measured']} of
        {coverage['scenarios_designed']}. Appendix N has the full register — three
        of the eight risks occurred, including the model being withdrawn
        mid-project.
    """)

    # ---------------------------------------------------------------- 13 artefact
    slide = _blank(prs)
    _heading(slide, "The artefact", "A working planner, not only an experiment")
    _stat_row(slide, [
        (f"{code['areas']['trip_planner']['lines']:,}", "lines in the system"),
        (f"{tests}", "tests, all passing"),
        (f"{measured.mcp_schema_stats()['tools_total']}", "schema-validated tools"),
        ("2", "interfaces — terminal\nand browser"),
    ], colour=ACCENT)
    _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
          [("It refuses a trip it cannot afford, and says what would work.",
            20, True, INK),
           ("Where a route has been recorded it uses the fare really quoted; "
            "where it has not, it says the figure is estimated.", 16, False,
            SOFT)], spacing=1.25)
    _notes(slide, """
        The feasibility gate is the part a user would notice. Ask for six nights in
        London on three thousand dollars and it declines, names the shortfall, and
        offers a shorter trip or a nearer city. It also distinguishes a measured
        price from an estimated one, because the price table was 52% below a real
        fare on the one route where both were known.
    """)

    # ---------------------------------------------------------------- 14 reflect
    slide = _blank(prs)
    _heading(slide, "What I would do differently",
             "The ordering mistake, and what it cost")
    _text(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.6),
          [("Measurement last was the wrong order.", 24, True, ACCENT),
           ("", 10, False, INK),
           ("Six agents became three because instrumentation showed where the "
            "requests went. That evidence arrived after the architecture was "
            "built, so the pivot cost a rebuild.", 17, False, INK),
           ("", 8, False, INK),
           ("Build the counter first, then the thing being counted.", 19, True,
            INK)], spacing=1.3)
    _notes(slide, """
        Honest answer to the likely question. The tool server came first because
        nothing works without it, then the message layer, then the agents, and the
        measurement infrastructure last. That ordering is why the first end-to-end
        run silently produced an ungrounded itinerary and nothing noticed. Reversed,
        the six-to-three pivot would have been a design decision rather than a
        rewrite.
    """)

    # ---------------------------------------------------------------- 15 close
    slide = _blank(prs)
    _fill(slide, DARK)
    _text(slide, Inches(1.1), Inches(1.6), Inches(11.2), Inches(0.5),
          [("CONCLUSION", 14, True, RGBColor(0x9C, 0xA3, 0xF5))])
    _text(slide, Inches(1.1), Inches(2.1), Inches(11.2), Inches(2.4),
          [("A model earns its cost where a step", 32, True, PAPER),
           ("needs a judgement you cannot write as code.", 32, True,
            RGBColor(0x6E, 0xE7, 0xDF)),
           ("Everywhere else it is a bill.", 32, True, PAPER)], spacing=1.15)
    _rule(slide, Inches(1.1), Inches(4.8), Inches(1.6), TEAL)
    _text(slide, Inches(1.1), Inches(5.15), Inches(11.2), Inches(1.6),
          [(f"{tokens('B') / tokens('D'):.0f}x fewer tokens, no loss of "
            f"groundedness, on {coverage['scenarios_measured']} scenario measured "
            f"{coverage['repeats_per_arm']} times.", 18, False,
            RGBColor(0xC7, 0xD2, 0xFE)),
           ("", 8, False, PAPER),
           ("Breadth is the gap. Everything needed to close it is committed and "
            "runs without credentials.", 15, False,
            RGBColor(0x9C, 0xA3, 0xF5))], spacing=1.3)
    _notes(slide, """
        Land on the claim and its limit in the same breath. The effect is large and
        the sample is narrow, and both are stated in the abstract. The cache, the
        harness and the scenarios are all committed, so the next person — or I,
        next month, when the quota resets — can widen it without rebuilding
        anything.
    """)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {os.path.relpath(path, ROOT)}")
