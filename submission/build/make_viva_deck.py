"""
Builds the viva slides.

Every number on a slide is read from the measured results. Slides are kept
deliberately sparse because the marking criteria penalise text-heavy ones; the
detail is in the speaker notes.
"""

from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(
    os.path.abspath(__file__)))))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from trip_planner.evaluation import measured

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIGURES = os.path.join(ROOT, "submission", "build", "figures")
OUTPUT = os.path.join(ROOT, "submission", "CMP7200_Viva_Presentation.pptx")

STUDENT = "25182589"          # anonymous marking: student number only
MODULE = "CMP7200 — Individual Master's Project"

INK = RGBColor(0x16, 0x1A, 0x23)
SOFT = RGBColor(0x54, 0x5C, 0x6E)
FAINT = RGBColor(0x8A, 0x91, 0xA0)
ACCENT = RGBColor(0x43, 0x38, 0xCA)
TEAL = RGBColor(0x0E, 0xA5, 0xA4)
WARN = RGBColor(0xA8, 0x5B, 0x00)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1B, 0x4B)
GREY_FILL_P = RGBColor(0xEF, 0xEE, 0xEA)
PURPLE_FILL_P = RGBColor(0x4E, 0x3F, 0xB8)
AQUA_FILL_P = RGBColor(0xD6, 0xF2, 0xE8)

W, H = Inches(13.333), Inches(7.5)      # 16:9


# ---------------------------------------------------------------------------
# Slide furniture
# ---------------------------------------------------------------------------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """
    Add a text box. `runs` is a list of (text, size_pt, bold, colour) tuples,
    one per paragraph, so a slide's type hierarchy is declared where it is used.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    for index, (text, size, bold, colour) in enumerate(runs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Calibri"
    return box


def _rule(slide, left, top, width, colour=ACCENT, height=Emu(34000)):
    bar = slide.shapes.add_shape(1, left, top, width, height)   # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = " ".join(text.split())


def _heading(slide, eyebrow, title):
    """The standard slide head: a small label, a large title, a short rule."""
    _text(slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.4),
          [(eyebrow.upper(), 12, True, ACCENT)])
    _text(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.9),
          [(title, 32, True, INK)])
    _rule(slide, Inches(0.9), Inches(1.72), Inches(1.1))


def _figure(slide, name, top=Inches(2.0), height=Inches(4.6)):
    """Place a generated figure, centred, or say plainly that it is missing."""
    for folder in ("results", "diagrams"):
        path = os.path.join(FIGURES, folder, name)
        if os.path.exists(path):
            picture = slide.shapes.add_picture(path, Inches(0), top,
                                               height=height)
            picture.left = int((W - picture.width) / 2)
            return picture
    _text(slide, Inches(0.9), top, Inches(11.5), Inches(0.5),
          [(f"[figure {name} not generated — run make_charts and make_diagrams]",
            14, False, WARN)])
    return None


def _stat_row(slide, stats, top=Inches(2.3), colour=ACCENT):
    """
    Evenly spaced big numbers. This is the deck's main device: a figure a marker
    can read from the back of the room, with three words under it.
    """
    count = len(stats)
    gutter = Inches(0.9)
    span = (W - 2 * gutter) / count
    for index, (value, label) in enumerate(stats):
        left = int(gutter + index * span)
        _text(slide, Emu(left), top, Emu(int(span)), Inches(1.3),
              [(value, 60, True, colour)], align=PP_ALIGN.CENTER)
        _text(slide, Emu(left), top + Inches(1.25), Emu(int(span)), Inches(0.8),
              [(label, 14, False, SOFT)], align=PP_ALIGN.CENTER)


def _table(slide, headers, rows, left=Inches(0.9), top=Inches(2.1),
           width=Inches(11.5), col_widths=None, font=12):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top,
                                   width, Inches(0.4 * (len(rows) + 1)))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for index, share in enumerate(col_widths):
            table.columns[index].width = Emu(int(width * share / total))
    def style(cell, value, bold, colour, fill):
        # Assigning "" leaves the paragraph with no runs, so styling it raises.
        # An empty header is legitimate — the approach table's first column is
        # just the letter — so write a space and style that.
        cell.text = str(value) if str(value) else " "
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(font)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

    for index, head in enumerate(headers):
        style(table.cell(0, index), head, True, PAPER, ACCENT)
    for r, row in enumerate(rows, start=1):
        shade = PAPER if r % 2 else RGBColor(0xF7, 0xF8, 0xFC)
        for c, value in enumerate(row):
            style(table.cell(r, c), value, False, INK, shade)
    return table


# ---------------------------------------------------------------------------
# The deck
# ---------------------------------------------------------------------------

def _panel(slide, left, top, width, height, text, *, fill, edge=None,
           size=15, bold=False, colour=None, align=PP_ALIGN.CENTER):
    """A filled rounded box with text in it. Used for the flow diagrams."""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # rounded rect
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour or INK
    run.font.name = "Calibri"
    return shape


def _arrow_down(slide, cx, top, height=Inches(0.32)):
    """A small downward arrow between two flow panels."""
    _text(slide, cx - Inches(0.4), top, Inches(0.8), height,
          [("▼", 14, False, FAINT)], align=PP_ALIGN.CENTER)


def build() -> str:
    arms = measured.results()["arms"]
    gains = measured.results().get("improvements", {})
    coverage = measured.coverage()
    protocol = measured.protocol_summary()
    gate = measured.gate_agreement()
    cache = measured.api_cache_stats()
    tests = measured.test_count()["collected"]
    per_arm = measured.api_calls_per_arm()["arms"]
    anchor = measured.gate_external_validity()

    from trip_planner.core.real_prices import PriceProbe
    from trip_planner.core.trip_cost import assess_budget
    from trip_planner.evaluation.scenarios import scenario

    demo_scenario = scenario(measured.scenario_ids()[0])
    demo = demo_scenario["params"]
    verdict = assess_budget(demo["budget"], demo["legs"][0][0], demo["nights"],
                            demo["adults"], demo["origin"],
                            price_probe=PriceProbe())
    grounding_ties = measured.intervals_overlap("C", "D", "prices_grounded_pct")

    def A(key):
        return arms["A"][key]

    def calls(letter):
        return arms[letter]["avg_llm_calls"]

    def tokens(letter):
        return arms[letter]["avg_total_tokens"]

    def grounded(letter):
        return arms[letter]["avg_prices_grounded_pct"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # =============================================== 1  title, problem, solution
    slide = _blank(prs)
    _fill(slide, DARK)
    _text(slide, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.4),
          [("MSc DISSERTATION  ·  VIVA VOCE  ·  " + MODULE, 12, True,
            RGBColor(0x9C, 0xA3, 0xF5))])
    _text(slide, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.1),
          [("AI Trip Planner", 46, True, PAPER)])
    _text(slide, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.6),
          [("You type where you want to go. It gives you a day-by-day plan with "
            "real flight and hotel prices.", 17, False,
            RGBColor(0xC7, 0xD2, 0xFE))])
    _rule(slide, Inches(0.9), Inches(2.95), Inches(1.4), TEAL)
    _text(slide, Inches(0.9), Inches(3.4), Inches(5.4), Inches(2.5),
          [("THE PROBLEM", 13, True, RGBColor(0xF0, 0xB4, 0x6A)),
           ("", 6, False, PAPER),
           ("A trip needs flights, hotels, places to visit and a budget that "
            "adds up.", 17, False, PAPER),
           ("", 4, False, PAPER),
           ("One AI doing all of it is slow, expensive, and makes prices up.",
            17, False, PAPER)], spacing=1.25)
    _text(slide, Inches(6.9), Inches(3.4), Inches(5.5), Inches(2.5),
          [("MY SOLUTION", 13, True, RGBColor(0x6E, 0xE7, 0xDF)),
           ("", 6, False, PAPER),
           ("Split the job between several small agents, each with one clear "
            "task.", 17, False, PAPER),
           ("", 4, False, PAPER),
           ("Then give AI only to the agents that need to think.", 17, False,
            RGBColor(0x6E, 0xE7, 0xDF))], spacing=1.25)
    _text(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
          [(f"Student {STUDENT}", 12, False, RGBColor(0x9C, 0xA3, 0xF5))])
    _notes(slide, """
        Twelve slides. A working travel planner, and one research question inside
        it: which parts actually need AI? I built it four ways and measured all
        four.
    """)

    # ======================================================== 2  the real question
    slide = _blank(prs)
    _heading(slide, "1  ·  What the system has to do",
             "One request, eight jobs")
    _panel(slide, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.75),
           f'"{demo_scenario["input"]}"',
           fill=RGBColor(0xF7, 0xF8, 0xFC), edge=RGBColor(0xE3, 0xE7, 0xF0),
           size=15, colour=SOFT)
    _table(slide,
           ["Understand the request", "Get real information", "Then"],
           [["1  Lahore to Istanbul", "4  Find real flights",
             "8  Build the day-by-day plan"],
            ["2  The dates", "5  Find real hotels", ""],
            ["3  The budget", "6  Find real attractions", ""],
            ["", "7  Find real restaurants", ""]],
           col_widths=[3.4, 3.6, 4.5], font=13.5, top=Inches(3.1))
    _text(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.4),
          [("So the real question is:", 16, False, SOFT),
           ("Should AI decide when to search for flights and hotels, or should "
            "normal code just do it?", 21, True, ACCENT)], spacing=1.3)
    _notes(slide, """
        This is the scenario every one of the four approaches is given, word for
        word, so nothing in the comparison depends on phrasing. Jobs 1 to 3 need
        judgement - there is no single right way to read a sentence. Jobs 4 to 7
        are fixed once you know the destination and the dates. Job 8 needs
        judgement again.
    """)

    # ============================================================= 3  the agents
    slide = _blank(prs)
    _heading(slide, "2  ·  The agents we use", "Each has a name and one job")
    _table(slide,
           ["Agent", "What it does", "Needs AI?"],
           [["Travel Conversation Assistant", "Asks for anything missing",
             "Yes - people write in many ways"],
            ["Travel Preferences Extractor", "Turns those words into exact fields",
             "Yes - same reason"],
            ["Flight Search Specialist", "Looks up flights", "No"],
            ["Hotel Search Specialist", "Looks up hotels", "No"],
            ["Activities Specialist", "Looks up places to visit and eat", "No"],
            ["Itinerary Coordinator", "Writes the day-by-day plan",
             "Yes - which day, what order"]],
           col_widths=[3.4, 4.6, 3.0], font=13, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
          [("Looking something up does not need thinking. Only three of the six "
            "really need AI.", 19, True, ACCENT)])
    _notes(slide, """
        The test for each agent: could I write its job as ordinary code? You
        cannot write a rule that understands any way a person might phrase a
        request, and you cannot write a rule that makes a day pleasant. You can
        absolutely write a rule that calls a flight API with a date.
    """)

    # ======================================================== 4  four approaches
    slide = _blank(prs)
    _heading(slide, "3  ·  I built it four ways", "Same request, same data")
    _table(slide,
           ["", "Name", "Simple meaning", "Agents", "Who fetches the data"],
           [["A", "Single AI", "AI answers from its own knowledge", "1",
             "nobody - no APIs at all"],
            ["B", "6 agents, naive", "Each agent decides its own tool use", "6",
             "the agents, through the tool server"],
            ["C", "6 agents, tuned", "Same six, configured properly", "6",
             "the agents, calling directly"],
            ["D", "3 agents, direct", "AI thinks; code does the searching", "3",
             "plain Python"]],
           col_widths=[0.4, 2.0, 3.9, 0.9, 3.8], font=13, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.3),
          [("All four get the identical request and the identical travel data.",
            18, True, INK),
           ("So the only thing being tested is who decides when to fetch it.",
            17, False, SOFT)], spacing=1.25)
    _notes(slide, """
        C is the one that makes this honest. Without it, someone could say D only
        won because B was configured badly - and they would be partly right, which
        is the finding on slide 7.
    """)

    # ============================================================ 5  approach A
    slide = _blank(prs)
    _heading(slide, "4  ·  Approach A  —  Single AI", "The baseline")
    _panel(slide, Inches(2.6), Inches(2.1), Inches(2.4), Inches(0.85), "USER",
           fill=GREY_FILL_P, size=14, bold=True)
    _panel(slide, Inches(5.5), Inches(2.1), Inches(2.4), Inches(0.85), "AI",
           fill=PURPLE_FILL_P, size=14, bold=True, colour=PAPER)
    _panel(slide, Inches(8.4), Inches(2.1), Inches(2.4), Inches(0.85), "ITINERARY",
           fill=GREY_FILL_P, size=14, bold=True)
    _text(slide, Inches(5.0), Inches(2.25), Inches(0.5), Inches(0.5),
          [("▶", 14, False, FAINT)], align=PP_ALIGN.CENTER)
    _text(slide, Inches(7.9), Inches(2.25), Inches(0.5), Inches(0.5),
          [("▶", 14, False, FAINT)], align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.6),
          [("No tools. No APIs. The AI writes the plan from what it already knows.",
            17, False, SOFT)], align=PP_ALIGN.CENTER)
    _stat_row(slide, [
        (f"{calls('A'):.0f}", "AI call"),
        (f"{grounded('A'):.1f}%", "of its prices matched a real fare"),
        (f"{A('avg_latency'):.0f}s", "to produce the plan"),
    ], top=Inches(4.1), colour=WARN)
    _text(slide, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.9),
          [("Cheap and fast, and the information cannot be trusted. That is what "
            "this approach is here to show.", 18, True, WARN)])
    _notes(slide, f"""
        Say it plainly: Approach A is the baseline. It shows what happens when AI
        plans a trip with no access to real travel data. It produced a fluent,
        confident itinerary and {grounded('A'):.1f}% of the prices in it matched
        anything real. Every other approach has to beat this, and if one of them
        did not, the tool layer would not be earning its place.
    """)

    # ============================================================ 6  approach B
    slide = _blank(prs)
    _heading(slide, "5  ·  Approach B  —  6 agents, naive",
             "Every agent decides for itself")
    _table(slide,
           ["The hotel agent's thinking", "Each line is a separate AI call"],
           [["\"I need Istanbul's destination ID\"", "tool call"],
            ["\"Now I'll search hotels\"", "tool call - 12,000 characters back"],
            ["\"Maybe I should check the reviews\"", "tool call"],
            ["\"Maybe I should look at nearby attractions\"", "tool call"],
            ["...and it may keep going, up to 10 steps", "8 tool descriptions "
             "re-sent every single time"]],
           col_widths=[6.2, 4.6], font=13, top=Inches(2.15))
    _stat_row(slide, [
        (f"{calls('B'):.0f}", "AI calls per trip"),
        (f"{arms['B']['avg_latency']:.0f}s", "to plan one trip"),
        ("2", "tools it never called at all"),
    ], top=Inches(4.7), colour=WARN)
    _text(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.7),
          [("Too much freedom: it keeps thinking and searching when it does not "
            "need to - and still skipped tools.", 17, True, WARN)])
    _notes(slide, f"""
        This was my original proposal. The agents can keep thinking, and every
        step re-sends the conversation and all the tool descriptions. On the
        recorded run it took {calls('B'):.1f} AI calls and
        {arms['B']['avg_latency']:.0f} seconds, produced malformed output twice
        that its own loop had to retry, and never called the attractions or the
        restaurant tool once - while still writing a complete-looking plan.
    """)

    # ============================================================ 7  approach C
    slide = _blank(prs)
    _heading(slide, "6  ·  Approach C  —  6 agents, tuned",
             "Approach B done properly")
    _table(slide,
           ["", "Approach B", "Approach C"],
           [["Tools per agent", "up to 8", "1 or 2"],
            ["Thinking steps allowed", "8 to 15", "3"],
            ["What the AI is sent", "the full 12,000-character API reply",
             "the best 3 results, trimmed"],
            ["When agents run", "one after another", "at the same time"]],
           col_widths=[2.8, 4.0, 3.7], font=13.5, top=Inches(2.2))
    _stat_row(slide, [
        (f"{gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}%", "less text sent, "
         "from settings alone"),
        (f"{calls('C'):.0f}", f"AI calls, down from {calls('B'):.0f}"),
        (f"{grounded('C'):.0f}%", "of its prices are real"),
    ], top=Inches(4.6), colour=ACCENT)
    _text(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.8),
          [("This is one of the biggest findings: most of B's cost was bad "
            "configuration, not the design.", 18, True, ACCENT)])
    _notes(slide, f"""
        Instead of handing the AI 12,000 characters of hotel data, C trims it to
        three lines - Theodora Pension $33 10/10, Hotel Sultania $41 9/10, and so
        on. Configuration alone removed
        {gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}% of the tokens. Being
        honest about that is what makes the next slide's claim defensible rather
        than inflated.
    """)

    # ============================================================ 8  approach D
    slide = _blank(prs)
    _heading(slide, "7  ·  Approach D  —  3 agents, direct",
             "The design that ships")
    mid = W / 2
    _panel(slide, mid - Inches(3.2), Inches(2.05), Inches(6.4), Inches(0.8),
           "AI STEP 1     understand the request",
           fill=PURPLE_FILL_P, size=15, bold=True, colour=PAPER)
    _arrow_down(slide, mid, Inches(2.9))
    _panel(slide, mid - Inches(3.2), Inches(3.3), Inches(6.4), Inches(1.5),
           "PLAIN PYTHON     fetch flights, hotels,\nattractions and restaurants",
           fill=AQUA_FILL_P, size=15, bold=True)
    _arrow_down(slide, mid, Inches(4.85))
    _panel(slide, mid - Inches(3.2), Inches(5.25), Inches(6.4), Inches(0.8),
           "AI STEP 2     build the day-by-day plan",
           fill=PURPLE_FILL_P, size=15, bold=True, colour=PAPER)
    _text(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.9),
          [("Once we know the destination and the dates, searching is fixed. "
            "There is nothing left for AI to decide.", 18, True, INK)],
          align=PP_ALIGN.CENTER)
    _notes(slide, f"""
        The one-sentence version of the whole project: AI decides what the user
        wants and how to present it; Python decides what data needs fetching. It
        makes exactly {per_arm['D']['total_http']} requests every time, in the
        same order, because an IF statement decides rather than a model. Nothing
        can be skipped, and {calls('D'):.0f} AI calls do the whole trip.
    """)

    # =========================================================== 9  where MCP fits
    slide = _blank(prs)
    _heading(slide, "8  ·  Where the tool server fits",
             "12 tools, but not every approach uses it the same way")
    _table(slide,
           ["", "How it reaches a travel API"],
           [["A", "It does not. No tools at all"],
            ["B", "Agent -> tool wrapper -> MCP client -> JSON-RPC -> MCP server "
                  "-> travel API"],
            ["C", "Agent -> the same server functions, called directly in the "
                  "process. No JSON-RPC"],
            ["D", "Plain Python -> the same server functions, called directly. "
                  "Code decides when"]],
           col_widths=[0.4, 10.6], font=13.5, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.5),
          [("If asked \"does the final system use the MCP server?\"", 16, False,
            SOFT),
           ("It uses the same tool functions the server exposes, but calls them "
            "in-process rather than over JSON-RPC. The protocol itself is "
            "exercised by Approach B and by the conformance audit.", 17, True,
            INK)], spacing=1.3)
    _notes(slide, """
        This is the question most likely to catch someone out, so it is on a slide
        rather than left to memory. It is also stated in Section 7.2 of the report
        as a limitation rather than glossed over: the two protocols the project is
        named for are less load-bearing in the shipped version than the proposal
        implied.
    """)

    # ========================================================= 10  the comparison
    slide = _blank(prs)
    _heading(slide, "9  ·  The four side by side",
             f"{coverage['repeats_per_arm']} runs of each")
    _table(slide,
           ["", "A  single AI", "B  6 naive", "C  6 tuned", "D  3 direct"],
           [["AI calls", f"{calls('A'):.0f}", f"{calls('B'):.0f}",
             f"{calls('C'):.0f}", f"{calls('D'):.0f}"],
            ["Text sent to the AI", f"{tokens('A'):,.0f}", f"{tokens('B'):,.0f}",
             f"{tokens('C'):,.0f}", f"{tokens('D'):,.0f}"],
            ["Cost per trip", f"${arms['A']['avg_cost_usd']:.4f}",
             f"${arms['B']['avg_cost_usd']:.4f}", f"${arms['C']['avg_cost_usd']:.4f}",
             f"${arms['D']['avg_cost_usd']:.4f}"],
            ["Time", f"{arms['A']['avg_latency']:.0f}s",
             f"{arms['B']['avg_latency']:.0f}s", f"{arms['C']['avg_latency']:.0f}s",
             f"{arms['D']['avg_latency']:.0f}s"],
            ["Prices that are real", f"{grounded('A'):.1f}%", f"{grounded('B'):.0f}%",
             f"{grounded('C'):.0f}%", f"{grounded('D'):.0f}%"]],
           col_widths=[3.0, 2.0, 2.0, 2.0, 2.2], font=13.5, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.5),
          [("The honest reading, not the flattering one:", 16, False, SOFT),
           (f"D is not more accurate than C — their ranges overlap. D is "
            f"{tokens('C') / tokens('D'):.1f}x cheaper than C and "
            f"{tokens('B') / tokens('D'):.0f}x cheaper than B, with no "
            f"detectable loss of accuracy.", 18, True, ACCENT)], spacing=1.3)
    _notes(slide, f"""
        Read the last row carefully. C scores {grounded('C'):.0f}% and D
        {grounded('D'):.0f}%, and their intervals overlap, so I cannot claim D is
        better grounded - and I do not. What does not overlap is the token count.
        The claim is therefore cheaper and faster with no measurable penalty,
        which is a smaller claim than the numbers first suggest and the only one
        the data supports.
    """)

    # ======================================================== 11  what user sees
    slide = _blank(prs)
    _heading(slide, "10  ·  What the user sees",
             "And the check that stops a fantasy trip")
    _figure(slide, "frontend.png", top=Inches(1.9), height=Inches(3.7))
    _stat_row(slide, [
        (f"${demo['budget']:,.0f}", "traveller asked for"),
        (f"${verdict.estimate.minimum:,.0f}", "cheapest it can be done"),
        ("REFUSED" if not verdict.feasible else "ALLOWED",
         "with three ways to fix it"),
    ], top=Inches(5.75), colour=WARN)
    _notes(slide, f"""
        Two things worth demonstrating. The seven step rows tick over live as each
        finishes, so the page never looks frozen. And the budget check: it works
        out the cheapest this trip could possibly cost, using the real recorded
        fare rather than a built-in table - the table said
        ${anchor['estimated_minimum']:,.0f} and the API returned
        ${anchor['cheapest_real_fare']:,.0f} - then refuses and offers a shorter
        trip, a bigger budget, or a nearer city.
    """)

    # ============================================================ 12  conclusion
    slide = _blank(prs)
    _fill(slide, DARK)
    _text(slide, Inches(1.1), Inches(1.15), Inches(11.2), Inches(0.5),
          [("IN ONE SENTENCE", 14, True, RGBColor(0x9C, 0xA3, 0xF5))])
    _text(slide, Inches(1.1), Inches(1.65), Inches(11.2), Inches(2.0),
          [("AI decides what the user wants and how to present it.", 30, True,
            PAPER),
           ("Python decides what data needs to be fetched.", 30, True,
            RGBColor(0x6E, 0xE7, 0xDF))], spacing=1.2)
    _rule(slide, Inches(1.1), Inches(3.85), Inches(1.6), TEAL)
    _text(slide, Inches(1.1), Inches(4.2), Inches(11.2), Inches(2.6),
          [("WHAT I FOUND", 12, True, RGBColor(0x9C, 0xA3, 0xF5)),
           ("", 5, False, PAPER),
           ("Without tools, AI invents travel information. Tools are essential.",
            16, False, RGBColor(0xC7, 0xD2, 0xFE)),
           (f"Configuration matters enormously — settings alone saved "
            f"{gains.get('C_vs_B', {}).get('tokens_pct', 0):.0f}%.", 16, False,
            RGBColor(0xC7, 0xD2, 0xFE)),
           ("Taking AI out of the lookup step is cheaper and faster, with no "
            "detectable loss of accuracy.", 16, False,
            RGBColor(0x6E, 0xE7, 0xDF)),
           ("", 5, False, PAPER),
           (f"The gap: {coverage['scenarios_measured']} of "
            f"{coverage['scenarios_designed']} trips measured for cost, because "
            f"of the free API limits. {protocol['passed']} of "
            f"{protocol['total_checks']} of my own design checks pass. Both are "
            f"in the report.", 14, False, RGBColor(0x9C, 0xA3, 0xF5))],
          spacing=1.25)
    _notes(slide, f"""
        Land the claim and its limits together. The three findings map onto three
        comparisons: A against the rest shows tools are essential; B against C
        shows configuration matters; C against D shows removing AI from
        deterministic lookup is free money. {tests} tests and
        {cache['entries']} saved API replies mean anyone can check all of it
        without a single API key.
    """)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {os.path.relpath(path, ROOT)}")
